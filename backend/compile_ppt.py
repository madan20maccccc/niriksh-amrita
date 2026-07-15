import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Define slide size (16:9 widescreen)
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Styling Constants
COLOR_PRIMARY = RGBColor(219, 39, 119)     # Amrita Pink (#DB2777)
COLOR_TEXT = RGBColor(30, 41, 59)          # Dark Slate (#1E293B)
COLOR_MUTED = RGBColor(100, 116, 139)      # Grey (#64748B)

# Image Paths
img_ocr = r"C:\Users\Madan M\.gemini\antigravity\brain\3078e600-92d3-4ff8-8f1e-739cdc1b9919\niriksh_vitals_ocr_mockup_1784088843679.png"
img_whatsapp = r"C:\Users\Madan M\.gemini\antigravity\brain\3078e600-92d3-4ff8-8f1e-739cdc1b9919\niriksh_whatsapp_alert_mockup_1784088860288.png"
img_sbar = r"C:\Users\Madan M\.gemini\antigravity\brain\3078e600-92d3-4ff8-8f1e-739cdc1b9919\niriksh_sbar_handover_mockup_1784088879765.png"

def apply_text_formatting(paragraph, font_name="Arial", size=13, color=COLOR_TEXT, bold=False):
    paragraph.font.name = font_name
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold

def create_slide_header(slide, title_text, category_text="NirikshAmrita Operational Design"):
    # Header box
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    
    # Category / Tag
    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    apply_text_formatting(p_cat, font_name="Arial", size=9, color=COLOR_PRIMARY, bold=True)
    
    # Main Slide Title
    p_title = tf.add_paragraph()
    p_title.text = title_text
    apply_text_formatting(p_title, font_name="Georgia", size=24, color=COLOR_TEXT, bold=True)
    p_title.space_before = Pt(3)

def add_bullets_to_textbox(tf, bullets):
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        if idx % 2 == 0:
            apply_text_formatting(p, font_name="Arial", size=14, color=COLOR_PRIMARY, bold=True)
            p.space_before = Pt(8) if idx > 0 else Pt(0)
        else:
            apply_text_formatting(p, font_name="Arial", size=11, color=COLOR_TEXT, bold=False)
            p.text = "• " + p.text

# Blank layout for custom layouts
slide_layout = prs.slide_layouts[6]

# ── SLIDE 1: TITLE SLIDE ───────────────────────────────────────
slide1 = prs.slides.add_slide(slide_layout)
tb1 = slide1.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.0))
tf1 = tb1.text_frame
tf1.word_wrap = True

p1 = tf1.paragraphs[0]
p1.text = "NirikshAmrita"
apply_text_formatting(p1, font_name="Georgia", size=56, color=COLOR_PRIMARY, bold=True)

p2 = tf1.add_paragraph()
p2.text = "Agentic Shift Handover & Clinical Early Warning Surveillance Platform"
apply_text_formatting(p2, font_name="Arial", size=20, color=COLOR_TEXT, bold=False)
p2.space_before = Pt(12)

p3 = tf1.add_paragraph()
p3.text = "Designed for Patient Safety and Escalation at Amrita Hospital, Faridabad"
apply_text_formatting(p3, font_name="Arial", size=14, color=COLOR_MUTED, bold=False)
p3.space_before = Pt(8)

p4 = tf1.add_paragraph()
p4.text = "Team Members: Madan.M, Anna Clara Mathew, K.L.C Aditya"
apply_text_formatting(p4, font_name="Arial", size=14, color=COLOR_TEXT, bold=True)
p4.space_before = Pt(45)


# ── SLIDE 2: CLINICAL CONTEXT & PROBLEM ────────────────────────
slide2 = prs.slides.add_slide(slide_layout)
create_slide_header(slide2, "Clinical Surveillance Limitations & Shift Handovers", "Clinical Problem Statement")
tb2 = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
tf2 = tb2.text_frame
tf2.word_wrap = True

bullets2 = [
    "Physiological Deterioration Latency",
    "In acute care settings, patients undergo gradual, multi-parameter physiological decline rather than sudden onset decompensation. Isolated point-in-time checks fail to capture progressive decline across shifts.",
    "Information Fragmentation During Rotations",
    "Ward handovers are statistically high-risk windows for information loss. Critical parameters (e.g. oxygen desaturation, rising glucose) are often poorly communicated or missed.",
    "Alarm Fatigue & Clinical Prioritization (Joint Commission Compliance)",
    "Standard clinical dashboards generate noisy, unprioritized alerts, causing cognitive overload for ward nurses. Systems must deliver explainable, prioritized warnings to prevent alarm desensitization."
]
add_bullets_to_textbox(tf2, bullets2)


# ── SLIDE 3: VISUAL FLOWCHART WORKFLOW ──────────────────────────
slide3 = prs.slides.add_slide(slide_layout)
create_slide_header(slide3, "NirikshAmrita Agentic Workflow Process Flowchart", "Clinical Safety Flowchart")

# Add 5 process rounded rectangles and arrows horizontally
steps_flow = [
    ("Step 1\nIngestion", "OCR camera scan reads bedside monitor digit values."),
    ("Step 2\nValidation", "Typo filtering agent checks parameters ranges."),
    ("Step 3\nNEWS2 Scoring", "EWS agent computes clinical risks score."),
    ("Step 4\nEscalation", "Bypasses client to alert Dr. Ramesh on WhatsApp."),
    ("Step 5\nHandover", "AI generates SBAR note translated in Malayalam.")
]

# Draw shapes programmatically
for idx, (title, desc) in enumerate(steps_flow):
    # Calculate horizontal layout positions
    left_pos = Inches(0.75 + idx * 2.4)
    top_pos = Inches(2.2)
    width_pos = Inches(1.9)
    height_pos = Inches(1.4)
    
    # Process box
    shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, width_pos, height_pos)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.color.rgb = COLOR_PRIMARY
    
    tf_flow = shape.text_frame
    tf_flow.word_wrap = True
    tf_flow.margin_top = Inches(0.1)
    tf_flow.margin_left = Inches(0.1)
    tf_flow.margin_right = Inches(0.1)
    
    p_title = tf_flow.paragraphs[0]
    p_title.text = title
    p_title.alignment = PP_ALIGN.CENTER
    apply_text_formatting(p_title, font_name="Arial", size=13, color=RGBColor(255, 255, 255), bold=True)
    
    # Text box for descriptions below the shape
    desc_tb = slide3.shapes.add_textbox(left_pos - Inches(0.15), Inches(3.8), width_pos + Inches(0.3), Inches(2.5))
    desc_tf = desc_tb.text_frame
    desc_tf.word_wrap = True
    desc_p = desc_tf.paragraphs[0]
    desc_p.text = desc
    desc_p.alignment = PP_ALIGN.CENTER
    apply_text_formatting(desc_p, font_name="Arial", size=11, color=COLOR_TEXT, bold=False)

    # Draw right-pointing arrow if not the last block
    if idx < 4:
        arrow = slide3.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, 
            left_pos + width_pos + Inches(0.1), 
            top_pos + Inches(0.55), 
            Inches(0.3), 
            Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_MUTED
        arrow.line.color.rgb = COLOR_MUTED


# ── SLIDE 4: WHY AGENTIC WORKFLOWS? ────────────────────────────
slide4 = prs.slides.add_slide(slide_layout)
create_slide_header(slide4, "The Paradigm Shift: Passive Dashboards vs. Agentic Surveillance", "Agentic Design Philosophy")
tb4 = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
tf4 = tb4.text_frame
tf4.word_wrap = True

bullets4 = [
    "Active Task Ownership & Closure Tracking",
    "Conventional systems rely entirely on user observation. An agentic system actively calculates early warning scores (NEWS2) and tracks alert acknowledgment until closure.",
    "Escalation SLAs & Resource Routing",
    "Integrates active escalation loops. If a nurse fails to acknowledge a critical flag within a set time (e.g., 15 minutes), the system escalates the warning to supervisors.",
    "Explainable Audit Trails",
    "Maintains legally auditable logs of response times, actions taken, and escalation paths, providing hospital quality assurance boards with complete visibility."
]
add_bullets_to_textbox(tf4, bullets4)


# ── SLIDE 5: TECHNICAL STACK OVERVIEW ──────────────────────────
slide5 = prs.slides.add_slide(slide_layout)
create_slide_header(slide5, "System Technical Stack & Core Integrations", "Technical Stack")
tb5 = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
tf5 = tb5.text_frame
tf5.word_wrap = True

bullets5 = [
    "Frontend Layer",
    "React 19, Vite, TanStack Start (for fast server-side hydration), TailwindCSS for custom glassmorphism styles, and Recharts for physiological trend lines.",
    "Backend & Ingestion APIs",
    "FastAPI (Python 3.10) for async concurrency, SQLite for transactional database records, and SQLAlchemy Object-Relational Mapping (ORM).",
    "AI Inference Layer",
    "Google Gemini API (using gemini-2.5-flash for vision parsing, handover summaries, and translation).",
    "Surveillance Protocols",
    "WebSockets for real-time telemetry broadcasts and audio alert synchronization."
]
add_bullets_to_textbox(tf5, bullets5)


# ── SLIDE 6: WHY THIS TECH STACK? ──────────────────────────────
slide6 = prs.slides.add_slide(slide_layout)
create_slide_header(slide6, "Technical Rationale & Architectural Decisions", "System Rationale")
tb6 = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
tf6 = tb6.text_frame
tf6.word_wrap = True

bullets6 = [
    "Asynchronous Processing (FastAPI)",
    "Enables high-concurrency ingestion of telemetry payloads without blocking computational agents. Async I/O prevents server lag when querying databases or external APIs.",
    "Optimistic State Synchronization (React & TanStack)",
    "Ensures clinical dashboards maintain absolute synchronicity with backend state changes, providing nurses with real-time updates without interface delays.",
    "Multimodal Cognition (Gemini 2.5 Flash)",
    "Selected for its low-latency image tokenization and contextual reasoning, powering real-time digitization of monitor screens and multilingual translation layers.",
    "ACID-Compliant Local Storage (SQLite & SQLAlchemy)",
    "Provides secure, transactional database records with low overhead, ensuring reliable local deployments during clinical simulations."
]
add_bullets_to_textbox(tf6, bullets6)


# ── SLIDE 7: STAGE 1 - CORE DATABASE & NEWS2 ───────────────────
slide7 = prs.slides.add_slide(slide_layout)
create_slide_header(slide7, "Stage 1: Core Database, Ingestion API & NEWS2 Core", "Stage 1 Core Ingestion")
tb7 = slide7.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
tf7 = tb7.text_frame
tf7.word_wrap = True

bullets7 = [
    "Clinical Data Models",
    "Define patient profiles with specific comorbidity flags (COPD, diabetes, hypertension) to customize warning thresholds.",
    "NEWS2 Calculator Agent",
    "Calculates clinical risk scores instantly upon vital entry. Dynamically adapts scoring ranges if the patient has a diagnosed COPD history.",
    "Physiological Range Validation",
    "Validation scripts filter out entry errors (e.g. BP entered as 800/60) before saving records, keeping database logs clean and accurate."
]
add_bullets_to_textbox(tf7, bullets7)


# ── SLIDE 8: STAGE 2 - GEMINI MULTIMODAL OCR ───────────────────
slide8 = prs.slides.add_slide(slide_layout)
create_slide_header(slide8, "Stage 2: Gemini Multimodal Telemetry OCR Ingestion", "Stage 2 Vision Ingestion")

tb8 = slide8.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
tf8 = tb8.text_frame
tf8.word_wrap = True

bullets8 = [
    "Vision OCR Ingestion",
    "Sends bedside monitor photos directly to gemini-2.5-flash to extract heart rate, blood pressure, SpO2, RR, and temperature values.",
    "Low-Cost Telemetry Parsing",
    "Bypasses complex physical hardware integration by letting nurses scan the monitor screen directly to populate the forms.",
    "Scanning UI Components",
    "Features a scanning line animation and a progress loader to show the ingestion status clearly."
]
add_bullets_to_textbox(tf8, bullets8)

if os.path.exists(img_ocr):
    slide8.shapes.add_picture(img_ocr, Inches(6.5), Inches(1.8), width=Inches(6.08), height=Inches(4.8))


# ── SLIDE 9: STAGE 3 - TRENDS & SBAR HANDOVER ──────────────────
slide9 = prs.slides.add_slide(slide_layout)
create_slide_header(slide9, "Stage 3: Cross-Shift Trend Reasoning & SBAR Note", "Stage 3 Clinical Reasoning")

tb9 = slide9.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
tf9 = tb9.text_frame
tf9.word_wrap = True

bullets9 = [
    "Shift-Wise Trend Calculations",
    "Analyzes the last 6 vital entries to identify gradual patient deterioration (e.g. SBP drops >10% or HR increases over 3 consecutive shifts).",
    "AI-Generated SBAR Summaries",
    "Processes patient history and trend data to write structured Situation, Background, Assessment, and Recommendation reports.",
    "Dynamic Line Graphs",
    "Renders historical parameter records (pulse, blood pressure, SpO2) in clean Recharts line charts."
]
add_bullets_to_textbox(tf9, bullets9)

if os.path.exists(img_sbar):
    slide9.shapes.add_picture(img_sbar, Inches(6.5), Inches(1.8), width=Inches(6.08), height=Inches(4.8))


# ── SLIDE 10: STAGE 4 - WHATSAPP & PHONE SIMULATOR ──────────────
slide10 = prs.slides.add_slide(slide_layout)
create_slide_header(slide10, "Stage 4: Outbound WhatsApp Alerts & Simulator", "Stage 4 Action & Alerting")

tb10 = slide10.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
tf10 = tb10.text_frame
tf10.word_wrap = True

bullets10 = [
    "Out-of-App Warning Gateway",
    "Sends WhatsApp alerts directly to the duty doctor's phone when a patient's NEWS2 score indicates orange or red risk.",
    "Actionable SLA Closure Protocols",
    "Alerts escalate to supervisors if the ward nurse does not acknowledge a red alert within 15 minutes.",
    "On-Screen Smartphone Mockup",
    "Features a floating iPhone mockup that simulates the WhatsApp chat and lock-screen alerts the doctor receives."
]
add_bullets_to_textbox(tf10, bullets10)

if os.path.exists(img_whatsapp):
    slide10.shapes.add_picture(img_whatsapp, Inches(6.5), Inches(1.8), width=Inches(6.08), height=Inches(4.8))


# ── SLIDE 11: STAGE 5 - WARD HEATMAP COMMAND CENTER ─────────────
slide11 = prs.slides.add_slide(slide_layout)
create_slide_header(slide11, "Stage 5: Ward Risk Heatmap & Command Center", "Stage 5 Oversight Command")
tb11 = slide11.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
tf11 = tb11.text_frame
tf11.word_wrap = True

bullets11 = [
    "Ward-Level Risk Heatmap",
    "Displays a grid of all beds colored by their active NEWS2 risk level (Red, Orange, Yellow, Green), giving supervisors a bird's-eye view.",
    "Interactive Tooltip Cards",
    "Hovering over a bed cell opens a small tooltip showing the patient's name, active warning description, and last actions taken.",
    "Surveillance Analytics Metrics",
    "Tracks average acknowledgment time, active alerts, occupancy logs, and rosters to monitor clinical workflow efficiency."
]
add_bullets_to_textbox(tf11, bullets11)


# ── SLIDE 12: EXTRA CLINICAL INNOVATIONS ───────────────────────
slide12 = prs.slides.add_slide(slide_layout)
create_slide_header(slide12, "Integrated System Innovations & Safety Policies", "Clinical Innovations")
tb12 = slide12.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
tf12 = tb12.text_frame
tf12.word_wrap = True

bullets12 = [
    "Malayalam & Telugu SBAR Translation",
    "Permits shift-handover summaries to be translated into regional languages in real-time, removing communication barriers.",
    "Clinical RAG Assistant Chatbot",
    "A patient-specific chatbot widget that allows nurses to query historical records, trends, and warning rationales.",
    "Deterioration Projections (AI Predictions)",
    "Uses regression modeling on vitals history to forecast boundaries breaches and alert nurses before they occur.",
    "Presentation Demonstration Mode",
    "A toggle switch that runs the OCR scanner, populates abnormal values, and triggers the WhatsApp overlay, simplifying live presentations."
]
add_bullets_to_textbox(tf12, bullets12)

# Save
output_path = r"D:\nursewatch-ai\NirikshAmrita_Project_Submission_V2.pptx"
prs.save(output_path)
print(f"Presentation with Flowchart saved successfully to: {output_path}")
