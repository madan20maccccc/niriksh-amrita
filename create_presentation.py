import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    # Color Palette (AIMS / NirikshAmrita theme: Rose/Pink/Violet/Deep Slate)
    DARK_BG = RGBColor(17, 24, 39)       # #111827 (Deep Slate)
    LIGHT_BG = RGBColor(253, 244, 245)   # #FDF4F5 (Soft Rose/White)
    PRIMARY_COLOR = RGBColor(219, 39, 119) # #DB2777 (Deep Pink / Rose)
    TEXT_DARK = RGBColor(31, 41, 55)      # #1F2937 (Charcoal)
    TEXT_LIGHT = RGBColor(243, 244, 246)  # #F3F4F6 (Off-white)
    MUTED_TEXT = RGBColor(107, 114, 128)  # #6B7280 (Gray)
    ACCENT_COLOR = RGBColor(124, 58, 237) # #7C3AED (Violet)

    # Helper function to style text
    def add_para(tf, text, font_size=Pt(14), color=TEXT_DARK, bold=False, italic=False, font_name="Segoe UI"):
        p = tf.add_paragraph() if tf.paragraphs and tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = font_size
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.space_after = Pt(8)
        return p

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # ---------------------------------------------
    # SLIDE 1: Title Slide (Dark Background)
    # ---------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, DARK_BG)

    # Decorative colored top border
    accent_bar = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = PRIMARY_COLOR
    accent_bar.line.fill.background()

    # Title Text Frame
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(4.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True

    add_para(tf1, "NirikshAmrita", font_size=Pt(54), color=PRIMARY_COLOR, bold=True)
    add_para(tf1, "AI Shift Safety & Early Warning System", font_size=Pt(28), color=TEXT_LIGHT, bold=True)
    add_para(tf1, "Intelligent Ward Surveillance and Automated Escalation Net", font_size=Pt(18), color=MUTED_TEXT, italic=True)
    
    # Team Box
    team_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.33), Inches(2.0))
    tf_team = team_box.text_frame
    tf_team.word_wrap = True
    add_para(tf_team, "Presented by Team NurseWatch-AI:", font_size=Pt(14), color=PRIMARY_COLOR, bold=True)
    add_para(tf_team, "• Madan.M (Project Lead)  • Anna Clara Mathew  • K.L.C Aditya", font_size=Pt(15), color=TEXT_LIGHT)
    add_para(tf_team, "School of AI, Amrita Vishwa Vidyapeetham", font_size=Pt(12), color=MUTED_TEXT)

    # ---------------------------------------------
    # SLIDE 2: Project Overview & Objectives (Light)
    # ---------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, LIGHT_BG)

    # Slide Header
    header_box = slide2.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.0))
    tf_h2 = header_box.text_frame
    add_para(tf_h2, "Project Overview & Objectives", font_size=Pt(36), color=PRIMARY_COLOR, bold=True)

    # Content Column Left
    content_box = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    tf_c2 = content_box.text_frame
    tf_c2.word_wrap = True

    add_para(tf_c2, "What is NirikshAmrita?", font_size=Pt(20), color=ACCENT_COLOR, bold=True)
    add_para(tf_c2, "NirikshAmrita is an agentic AI-driven hospital ward surveillance system built to monitor patient telemetry in real-time, predict medical deterioration using standard protocols, and streamline clinical workflows.", font_size=Pt(15), color=TEXT_DARK)
    
    add_para(tf_c2, "\nCore Project Objectives:", font_size=Pt(20), color=ACCENT_COLOR, bold=True)
    add_para(tf_c2, "• Dynamic NEWS2 Scoring: Auto-calculates early warning scores instantly from vitals data.", font_size=Pt(14), color=TEXT_DARK)
    add_para(tf_c2, "• Prevent Alarm Fatigue: Intelligently routes and escalates alerts only when action is missed.", font_size=Pt(14), color=TEXT_DARK)
    add_para(tf_c2, "• Nurse Workload Reduction: Automatically compiles shift handover summaries (SBAR reports).", font_size=Pt(14), color=TEXT_DARK)
    add_para(tf_c2, "• HIPAA Compliance: Anonymizes presentation telemetry and restricts private data behind authorization.", font_size=Pt(14), color=TEXT_DARK)

    # ---------------------------------------------
    # SLIDE 3: Technology Stack (Light)
    # ---------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, LIGHT_BG)

    header_box = slide3.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.0))
    tf_h3 = header_box.text_frame
    add_para(tf_h3, "The Technology Stack", font_size=Pt(36), color=PRIMARY_COLOR, bold=True)

    # Left Column (Frontend)
    col1_box = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(5.0))
    tf_col1 = col1_box.text_frame
    tf_col1.word_wrap = True
    add_para(tf_col1, "Frontend (Client Workspace)", font_size=Pt(20), color=ACCENT_COLOR, bold=True)
    add_para(tf_col1, "• React 19 & TypeScript: Modern component model.", font_size=Pt(14), color=TEXT_DARK)
    add_para(tf_col1, "• TanStack Start & Nitro: Server-side rendering (SSR) for instant layout hydration.", font_size=Pt(14), color=TEXT_DARK)
    add_para(tf_col1, "• Recharts: Renders hemodynamic, metabolic, thermal, and respiratory trends interactively.", font_size=Pt(14), color=TEXT_DARK)
    add_para(tf_col1, "• TailwindCSS: Premium custom color theme aligned with AIMS styling.", font_size=Pt(14), color=TEXT_DARK)

    # Right Column (Backend)
    col2_box = slide3.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(5.0))
    tf_col2 = col2_box.text_frame
    tf_col2.word_wrap = True
    add_para(tf_col2, "Backend & AI Layer", font_size=Pt(20), color=ACCENT_COLOR, bold=True)
    add_para(tf_col2, "• FastAPI (Python): Async backend servicing REST routes and WebSocket live broadcasts.", font_size=Pt(14), color=TEXT_DARK)
    add_para(tf_col2, "• SQLite & SQLAlchemy ORM: Embedded relational database for storage.", font_size=Pt(14), color=TEXT_DARK)
    add_para(tf_col2, "• Gemini API & RAG Agent: Powers the clinical chatbot and auto-generates SBAR summaries.", font_size=Pt(14), color=TEXT_DARK)
    add_para(tf_col2, "• Docker & Hugging Face: Deployed as a single container running both frontend and backend on port 7860.", font_size=Pt(14), color=TEXT_DARK)

    # ---------------------------------------------
    # SLIDE 4: Operational Workflow & Flowchart (Light)
    # ---------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, LIGHT_BG)

    header_box = slide4.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.0))
    tf_h4 = header_box.text_frame
    add_para(tf_h4, "Platform Operational Workflow", font_size=Pt(36), color=PRIMARY_COLOR, bold=True)

    flow_box = slide4.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    tf_flow = flow_box.text_frame
    tf_flow.word_wrap = True

    add_para(tf_flow, "1. Bedside Rounds Entry", font_size=Pt(18), color=PRIMARY_COLOR, bold=True)
    add_para(tf_flow, "Nurse inputs Systolic BP, HR, SpO2, and Temp. Shift is automatically detected based on system clock.", font_size=Pt(14), color=TEXT_DARK)

    add_para(tf_flow, "2. Asynchronous Multi-Agent Assessment", font_size=Pt(18), color=PRIMARY_COLOR, bold=True)
    add_para(tf_flow, "Validation Agent sanitizes inputs → EWS Agent calculates NEWS2 → Trend Agent runs regression predictions.", font_size=Pt(14), color=TEXT_DARK)

    add_para(tf_flow, "3. Broadcast & Safety Net Escalation", font_size=Pt(18), color=PRIMARY_COLOR, bold=True)
    add_para(tf_flow, "Vitals score updates all ward boards instantly. If high risk, a 4-tier cascade begins:\n"
                     "   [Level 1: Nurse Alert] → [Level 2: Supervisor Alert] → [Level 3: Doctor Page] → [Level 4: Closure Check]", font_size=Pt(14), color=TEXT_DARK)

    # ---------------------------------------------
    # SLIDE 5: Key Platform Features (Light)
    # ---------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, LIGHT_BG)

    header_box = slide5.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.0))
    tf_h5 = header_box.text_frame
    add_para(tf_h5, "Key Implemented Features", font_size=Pt(36), color=PRIMARY_COLOR, bold=True)

    features_box = slide5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    tf_f = features_box.text_frame
    tf_f.word_wrap = True

    add_para(tf_f, "✔ Dynamic Patient Registry & Nurse Shift Rostering", font_size=Pt(16), color=TEXT_DARK, bold=True)
    add_para(tf_f, "✔ Interactive Vital charts across 4 physiological domains", font_size=Pt(16), color=TEXT_DARK, bold=True)
    add_para(tf_f, "✔ Live Audio Alarms & WebSocket status synchronization", font_size=Pt(16), color=TEXT_DARK, bold=True)
    add_para(tf_f, "✔ Multilingual SBAR Handover: Real-time translation to Malayalam and Telugu", font_size=Pt(16), color=PRIMARY_COLOR, bold=True)
    add_para(tf_f, "✔ RAG-Based Clinical Chatbot: Talk to the patient's record to get instant trend explanations", font_size=Pt(16), color=TEXT_DARK, bold=True)
    add_para(tf_f, "✔ Safe Seeding & Containerized Start Scripts: Ready-to-go deployment pipeline", font_size=Pt(16), color=TEXT_DARK, bold=True)

    # ---------------------------------------------
    # SLIDE 6: Roadmap & Next Steps (Dark Background)
    # ---------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, DARK_BG)

    header_box = slide6.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.0))
    tf_h6 = header_box.text_frame
    add_para(tf_h6, "Future Roadmap & Next Steps", font_size=Pt(36), color=PRIMARY_COLOR, bold=True)

    roadmap_box = slide6.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
    tf_r = roadmap_box.text_frame
    tf_r.word_wrap = True

    add_para(tf_r, "Phase 1: Physical IoT Integrations", font_size=Pt(20), color=PRIMARY_COLOR, bold=True)
    add_para(tf_r, "Incorporate Bluetooth / Wi-Fi based pulse oximeters and blood pressure cuffs to eliminate manual entry errors entirely.", font_size=Pt(14), color=TEXT_LIGHT)

    add_para(tf_r, "Phase 2: Medical LLM Fine-Tuning", font_size=Pt(20), color=PRIMARY_COLOR, bold=True)
    add_para(tf_r, "Adapt fine-tuned medical models on local servers to provide automated guideline suggestions (e.g. sepsis protocols) immediately after a high NEWS2 triggers.", font_size=Pt(14), color=TEXT_LIGHT)

    add_para(tf_r, "Phase 3: Administrative Predictive Analytics", font_size=Pt(20), color=PRIMARY_COLOR, bold=True)
    add_para(tf_r, "Develop ML models to forecast ICU bed availability, hospital stay durations, and patient readmission probabilities.", font_size=Pt(14), color=TEXT_LIGHT)

    # Save presentation
    filename = "niriksh_amrita_presentation.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as {filename}")

if __name__ == "__main__":
    create_deck()
